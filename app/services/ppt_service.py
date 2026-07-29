from pptx import Presentation
from typing import Dict, Any, List


class PPTService:
    """PPT service with accurate slide-level metadata extraction."""

    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract all text from PPT."""
        try:
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"[PPTService] Error extracting text: {str(e)}")
            return ""

    @staticmethod
    def extract_text_with_metadata(file_path: str) -> Dict[str, Any]:
        """Extract text with slide numbers and metadata."""
        try:
            presentation = Presentation(file_path)
            slides = []
            full_text = ""

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                slides.append({
                    'page_num': slide_num,
                    'text': slide_text,
                })
                full_text += slide_text + "\n"

            doc_metadata = {}
            try:
                props = presentation.core_properties
                safe_attrs = ['title', 'author', 'subject', 'keywords', 'comments', 'category']
                for attr in safe_attrs:
                    try:
                        if hasattr(props, attr):
                            value = getattr(props, attr)
                            if value:
                                doc_metadata[attr] = str(value)
                    except:
                        pass
            except:
                pass

            return {
                'text': full_text,
                'pages': slides,
                'total_pages': len(slides),
                'metadata': doc_metadata
            }
        except Exception as e:
            print(f"[PPTService] Error extracting metadata: {str(e)}")
            return {'text': '', 'pages': [], 'total_pages': 0, 'metadata': {}}